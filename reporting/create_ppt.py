from __future__ import print_function
from pptx import Presentation
from pptx.util import Inches, Pt
import argparse
import numpy as np
from datetime import date
import jsonschema, json
import matplotlib.pyplot as plt
from summarize_metrics import summarize_metrics, BAAThresholds
from pathlib import Path
from argparse import ArgumentParser, ArgumentError
import pandas as pd
from operator import add


def directory_type(arg_string):
    """
    Allows arg parser to handle directories
    :param arg_string: A path, relative or absolute to a folder
    :return: A python pure path object to a directory.
    """
    directory_path = Path(arg_string)
    if directory_path.exists() and directory_path.is_dir():
        return str(directory_path)
    raise ArgumentError("{} is not a valid directory.".format(arg_string))


def list_type(arg_string):
    """
    Allows arg parser to handle strings
    :param arg_string: a python syntaxed list, as a string
    :return: A python list
    """
    arguments = arg_string.split()
    return arguments


def iter_cells(table):
    for row in table.rows:
        for cell in row.cells:
            yield cell


def df_to_table(slide, df, left, top, width, height, colnames=None):
    """Converts a Pandas DataFrame to a PowerPoint table on the given
    Slide of a PowerPoint presentation.
    The table is a standard Powerpoint table, and can easily be modified with the Powerpoint tools,
    for example: resizing columns, changing formatting etc.
    Arguments:
     - slide: slide object from the python-pptx library containing the slide on which you want the table to appear
     - df: Pandas DataFrame with the data
    Optional arguments:
     - colnames
     https://github.com/robintw/PandasToPowerpoint/blob/master/PandasToPowerpoint.py
     """
    rows, cols = df.shape
    res = slide.shapes.add_table(rows + 1, cols, left, top, width, height)

    if colnames is None:
        colnames = list(df.columns)

    # Insert the column names
    for col_index, col_name in enumerate(colnames):
        # Column names can be tuples
        if not isinstance(col_name, str):
            col_name = " ".join(col_name)
        res.table.cell(0, col_index).text = col_name

    m = df.values

    for row in range(rows):
        for col in range(cols):
            val = m[row, col]
            text = str(val)
            res.table.cell(row + 1, col).text = text

    for cell in iter_cells(res.table):
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(14)


def parse_args():
    """ Setup the input and output arguments for the script
    Return the parsed input and output files
    """
    parser = argparse.ArgumentParser(description='Create ppt report')
    parser.add_argument('-infile',
                        type=argparse.FileType('r'),
                        help='Powerpoint file used as the template',
                        required=True)
    parser.add_argument('-outfile',
                        type=argparse.FileType('w'),
                        help='Output powerpoint report file',
                        required=True)
    parser.add_argument('-rootdir',
                        type=directory_type,
                        help='Root directory for all teams',
                        required=True)
    parser.add_argument('-teams',
                        nargs='+',
                        help='Name of all teams meant to be evaluated',
                        required=True)
    parser.add_argument('-aois',
                        nargs='+',
                        help='All AOIs meant to be evaluated',
                        required=True)
    return parser.parse_args()


def create_team_based_aoi_scores_slide(baa_thresholds, prs, team_scores, classifications):
    if type(classifications) is int:
        classifications = [classifications]
    metrics_names = {'Metrics': ['2D Correctness', '2D Completeness', '2D IOU', '3D Correctness',
                                 '3D Completeness', '3D IOU', 'Geolocation Error', 'H-RMSE', 'Z-RMSE',
                                 'Mean Model Size (MB)', 'Textures (qualitative)', 'Model fitting (qualitative)',
                                 'Model fitting (qualitative)', 'Wall Clock Time (hrs/sq. km)',
                                 'Cost per sq. km ($USD)']}
    phase_1a_thresholds = {'Phase 1A Thresholds': [baa_thresholds.correctness_2d[0],
                                                   baa_thresholds.completeness_2d[0],
                                                   baa_thresholds.jaccard_index_2d[0],
                                                   baa_thresholds.correctness_3d[0],
                                                   baa_thresholds.completeness_3d[0],
                                                   baa_thresholds.jaccard_index_3d[0],
                                                   baa_thresholds.geolocation_error[0]]}
    phase_2b_thresholds = {'Phase 2B Thresholds': [baa_thresholds.correctness_2d[3],
                                                   baa_thresholds.completeness_2d[3],
                                                   baa_thresholds.jaccard_index_2d[3],
                                                   baa_thresholds.correctness_3d[3],
                                                   baa_thresholds.completeness_3d[3],
                                                   baa_thresholds.jaccard_index_3d[3],
                                                   baa_thresholds.geolocation_error[3]]}
    df_metrics_names = pd.DataFrame(data=metrics_names)
    df_1a_thresholds = pd.DataFrame(data=phase_1a_thresholds)
    df_2b_thresholds = pd.DataFrame(data=phase_2b_thresholds)
    metrics_aoi = {}
    for team in team_scores:
        team_aoi_metrics_slide = prs.slide_layouts[4]
        slide = prs.slides.add_slide(team_aoi_metrics_slide)
        title = slide.shapes.title
        title.text = (team + " - Site Scores - {0}").format(','.join(str(e) for e in classifications))

        for aoi in team_scores[team]:
            if aoi not in metrics_aoi.keys():
                metrics_aoi[aoi] = {}
            metrics_aoi[aoi].update(team_scores[team][aoi].results)

        df_team_metric = None
        for aoi in metrics_aoi:
            metrics_2d = [0, 0, 0]
            metrics_3d = [0, 0, 0]
            other_metrics = [0, 0, 0]
            for cls in classifications:
                metrics_2d = list(map(add, metrics_2d,
                                      [metrics_aoi[aoi]['threshold_geometry'][cls]['2D']['completeness'],
                                       metrics_aoi[aoi]['threshold_geometry'][cls]['2D']['correctness'],
                                       metrics_aoi[aoi]['threshold_geometry'][cls]['2D']['jaccardIndex']]))
                metrics_3d = list(map(add, metrics_3d,
                                      [metrics_aoi[aoi]['threshold_geometry'][cls]["3D"]['completeness'],
                                       metrics_aoi[aoi]['threshold_geometry'][cls]['3D']['correctness'],
                                       metrics_aoi[aoi]['threshold_geometry'][cls]['3D']['jaccardIndex']]))
                other_metrics = list(map(add, other_metrics,
                                         [metrics_aoi[aoi]["gelocation_error"],
                                          metrics_aoi[aoi]['relative_accuracy'][cls]["hrmse"],
                                          metrics_aoi[aoi]['relative_accuracy'][cls]["zrmse"]]))
            metrics_2d = [np.round(x / classifications.__len__(), decimals=2) for x in metrics_2d]
            metrics_3d = [np.round(x / classifications.__len__(), decimals=2) for x in metrics_3d]
            other_metrics = [np.round(x / classifications.__len__(), decimals=2) for x in other_metrics]
            metrics_column = {aoi: metrics_2d + metrics_3d + other_metrics}
            df_aoi_metric = pd.DataFrame(data=metrics_column)
            if df_team_metric is not None:
                df_team_metric = pd.concat([df_team_metric, df_aoi_metric], axis=1)
            else:
                df_team_metric = pd.DataFrame(data=metrics_column)

        df_full_table = pd.concat([df_metrics_names,
                                   df_1a_thresholds,
                                   df_2b_thresholds,
                                   df_team_metric],
                                  axis=1)
        top = Inches(1)
        left = Inches(0.5)
        width = Inches(12.3)
        height = Inches(5.6)
        df_to_table(slide, df_full_table, left, top, width, height)


def create_mean_scores_by_site_slide(baa_thresholds, prs, averaged_results, aois, classifications):
    # Check how many classifications we deal with
    if type(classifications) is int:
        classifications = [classifications]
    # Create Table of all teams
    # Create columns from summarized results
    df_performer_metrics = None
    metrics_names = {'Metrics': ['2D Correctness', '2D Completeness', '2D IOU', '3D Correctness',
                                 '3D Completeness', '3D IOU', 'Geolocation Error', 'H-RMSE', 'Z-RMSE',
                                 'Mean Model Size (MB)', 'Textures (qualitative)', 'Model fitting (qualitative)',
                                 'Model fitting (qualitative)', 'Wall Clock Time (hrs/sq. km)',
                                 'Cost per sq. km ($USD)']}
    phase_1a_thresholds = {'Phase 1A Thresholds': [baa_thresholds.correctness_2d[0],
                                                   baa_thresholds.completeness_2d[0],
                                                   baa_thresholds.jaccard_index_2d[0],
                                                   baa_thresholds.correctness_3d[0],
                                                   baa_thresholds.completeness_3d[0],
                                                   baa_thresholds.jaccard_index_3d[0],
                                                   baa_thresholds.geolocation_error[0]]}
    phase_2b_thresholds = {'Phase 2B Thresholds': [baa_thresholds.correctness_2d[3],
                                                   baa_thresholds.completeness_2d[3],
                                                   baa_thresholds.jaccard_index_2d[3],
                                                   baa_thresholds.correctness_3d[3],
                                                   baa_thresholds.completeness_3d[3],
                                                   baa_thresholds.jaccard_index_3d[3],
                                                   baa_thresholds.geolocation_error[3]]}
    df_metrics_names = pd.DataFrame(data=metrics_names)
    df_1a_thresholds = pd.DataFrame(data=phase_1a_thresholds)
    df_2b_thresholds = pd.DataFrame(data=phase_2b_thresholds)
    for team in averaged_results:
        metrics_2d = [0, 0, 0]
        metrics_3d = [0, 0, 0]
        other_metrics = [0, 0, 0]
        for cls in classifications:
            metrics_2d = list(map(add, metrics_2d, [averaged_results[team][cls]["2d_completeness"],
                                                    averaged_results[team][cls]["2d_correctness"],
                                                    averaged_results[team][cls]["2d_jaccard_index"]]))
            metrics_3d = list(map(add, metrics_3d, [averaged_results[team][cls]["3d_completeness"],
                                                    averaged_results[team][cls]["3d_correctness"],
                                                    averaged_results[team][cls]["3d_jaccard_index"]]))
            other_metrics = list(map(add, other_metrics, [averaged_results[team]["geolocation_error"],
                                                          averaged_results[team][cls]["hrmse"],
                                                          averaged_results[team][cls]["zrmse"]]))
        metrics_2d = [np.round(x / classifications.__len__(), decimals=2) for x in metrics_2d]
        metrics_3d = [np.round(x / classifications.__len__(), decimals=2) for x in metrics_3d]
        other_metrics = [np.round(x / classifications.__len__(), decimals=2) for x in other_metrics]
        metrics_column = {team: metrics_2d + metrics_3d + other_metrics}
        df_team_metrics = pd.DataFrame(data=metrics_column)
        if df_performer_metrics is not None:
            df_performer_metrics = pd.concat([df_performer_metrics, df_team_metrics], axis=1)
        else:
            df_performer_metrics = pd.DataFrame(data=metrics_column)

    df_full_table = pd.concat([df_metrics_names,
                               df_1a_thresholds,
                               df_2b_thresholds,
                               df_performer_metrics],
                              axis=1)

    summary_metrics_slide_layout = prs.slide_layouts[4]
    slide = prs.slides.add_slide(summary_metrics_slide_layout)
    title = slide.shapes.title
    title.text = "Mean Scores from {0} - Class: {1}".format('-'.join(aois), ','.join(str(e) for e in classifications))
    top = Inches(1)
    left = Inches(0.5)
    width = Inches(12.3)
    height = Inches(5.6)
    df_to_table(slide, df_full_table, left, top, width, height)


def create_metrics_images_slide(prs, aoi, configs):
    metrics_images_layout = prs.slide_layouts[4]
    slide = prs.slides.add_slide(metrics_images_layout)
    title = slide.shapes.title
    title.text = aoi
    for team in configs:
        # Get output results prefix
        search_path = configs[team][aoi]['path'].parent
        file_prefix = Path(configs[team][aoi]['INPUT.TEST']['DSMFilename']).name
        file_path_prefix = Path(search_path, file_prefix)
        suffixes = ["_000_objectwise_obj3dJaccardIndex.png", "_000_objectwise_objHRMSE.png",
                    "_000_relVertAcc_hgtErr_clipped.png"]
        for suffix in suffixes:
            filename_path = Path(str(file_path_prefix.absolute()) + suffix)
            if filename_path.is_file():
                # TODO: Figure out how to space images
                slide.shapes.add_picture(str(filename_path.absolute()), 0, 0)


def create_title_slide(prs):
    # Create a title slide first
    title_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "CORE3D Metrics Report"
    subtitle.text = "Generated on {:%m-%d-%Y}".format(date.today())
    subtitle_box = slide.placeholders[13]
    subtitle_box.text = "JHU/APL"


def create_ppt(layout_slides_input, output_slides, team_scores, averaged_results, baa_thresholds, aois, configs):
    """ Take the input powerpoint file and use it as the template for the output
    file.
    """
    prs = Presentation(layout_slides_input)
    # Use the output from analyze_ppt to understand which layouts and placeholders
    # to use
    create_title_slide(prs)

    create_mean_scores_by_site_slide(baa_thresholds, prs, averaged_results, aois, 6)
    create_mean_scores_by_site_slide(baa_thresholds, prs, averaged_results, aois, 17)
    create_mean_scores_by_site_slide(baa_thresholds, prs, averaged_results, aois, [6, 17])

    create_team_based_aoi_scores_slide(baa_thresholds, prs, team_scores, 6)
    create_team_based_aoi_scores_slide(baa_thresholds, prs, team_scores, 17)
    create_team_based_aoi_scores_slide(baa_thresholds, prs, team_scores, [6, 17])

    for aoi in aois:
        create_metrics_images_slide(prs, aoi, configs)

    prs.save(output_slides)


def main():
    args = parse_args()
    baa_thresholds = BAAThresholds()
    averaged_results, team_scores, configs = summarize_metrics(args.rootdir, args.teams, args.aois)
    create_ppt(args.infile.name, args.outfile.name, team_scores, averaged_results, baa_thresholds, args.aois, configs)


if __name__ == "__main__":
    main()
